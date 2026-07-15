from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "src" / "moodle" / "local_hubredirect" / "office_templates"


def create_docx(path: Path) -> None:
    document = Document()
    document.add_paragraph("")
    document.save(path)


def create_xlsx(path: Path) -> None:
    workbook = Workbook()
    workbook.active.title = "Sheet1"
    workbook.save(path)


def create_pptx(path: Path) -> None:
    deck = Presentation()
    blank = deck.slide_layouts[6]
    deck.slides.add_slide(blank)
    deck.save(path)


def create_pdf(path: Path) -> None:
    pdf = canvas.Canvas(str(path))
    pdf.showPage()
    pdf.save()


def main() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    create_docx(OUT / "blank.docx")
    create_xlsx(OUT / "blank.xlsx")
    create_pptx(OUT / "blank.pptx")
    create_pdf(OUT / "blank.pdf")


if __name__ == "__main__":
    main()
